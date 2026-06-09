"""Generate 3 Key Benefit slides for Shopify App Store at 1600x1200 (4:3).

Uses the Remaker Digital PPTX template, adjusts to 4:3 aspect ratio,
and creates 3 slides highlighting Agent Red's key differentiators.

Output: PPTX-template-and-skill/output/key-benefit-slides.pptx
The owner exports each slide as PNG at 1600x1200 from PowerPoint.
"""

import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# Paths
TEMPLATE = Path(__file__).resolve().parent.parent / "PPTX-template-and-skill" / "remaker-digital-slide-template.pptx"
OUTPUT = Path(__file__).resolve().parent.parent / "PPTX-template-and-skill" / "output" / "key-benefit-slides.pptx"

# Target dimensions: 1600x1200 pixels at 96 DPI = 16.667" x 12.5"
# In EMU (English Metric Units): 1 inch = 914400 EMU
TARGET_WIDTH = Emu(round(16.667 * 914400))  # ~15240480 EMU
TARGET_HEIGHT = Emu(round(12.5 * 914400))  # ~11430000 EMU

# Brand colors
BRAND_RED = RGBColor(0xFF, 0x36, 0x21)  # #ff3621
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_BG = RGBColor(0x0A, 0x0A, 0x0A)  # Chrome color from design system
SURFACE = RGBColor(0x1F, 0x1F, 0x1F)  # Surface color from design system

# 3 Key Benefits for Shopify App Store
# These are marketing-focused, merchant-facing value propositions
SLIDES = [
    {
        "title": "AI that remembers every customer",
        "subtitle": "Persistent Customer Memory",
        "points": [
            "Builds a profile from every interaction automatically",
            "Recalls purchase history, preferences, and context",
            "Delivers personalized responses from the first message",
            "4-layer memory: profile, history, patterns, dedicated AI",
        ],
        "callout": "No competitor offers per-customer AI memory",
    },
    {
        "title": "Instant answers, always on",
        "subtitle": "6-Agent AI pipeline",
        "points": [
            "Intent classification routes queries in milliseconds",
            "Knowledge base search with hybrid AI retrieval",
            "Real-time streaming responses via SSE",
            "Safety-first: every response validated before delivery",
        ],
        "callout": "< 2 second response time, 24/7/365",
    },
    {
        "title": "One dashboard, complete control",
        "subtitle": "Built for Shopify merchants",
        "points": [
            "Conversation inbox with full customer context",
            "Knowledge base with document upload and AI search",
            "Widget customizer with live preview",
            "Usage analytics, team management, and billing",
        ],
        "callout": "Starting at $149/mo - 4 to 21x cheaper than competitors",
    },
]


def scale_position(original_val, original_slide_dim, new_slide_dim):
    """Scale a position/size proportionally to new slide dimensions."""
    return int(original_val * new_slide_dim / original_slide_dim)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Montserrat",
):
    """Add a text box to the slide with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, bullets, font_size=20, color=WHITE, bullet_color=None):
    """Add a bulleted text list to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Add bullet character + text
        p.text = bullet_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Montserrat"
        p.space_after = Pt(12)

        # Set bullet
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buFont = etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buFont")
        buFont.set("typeface", "Arial")
        buChar = etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
        buChar.set("char", "\u2022")  # bullet character

        # Bullet color
        if bullet_color:
            buClr = etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buClr")
            srgbClr = etree.SubElement(buClr, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
            srgbClr.set("val", f"{bullet_color[0]:02X}{bullet_color[1]:02X}{bullet_color[2]:02X}")

        # Indent
        pPr.set("marL", str(Emu(Inches(0.4))))
        pPr.set("indent", str(Emu(Inches(-0.25))))

    return txBox


def add_accent_line(slide, left, top, width, color=BRAND_RED, thickness=Pt(4)):
    """Add a horizontal accent line."""
    from pptx.util import Emu as E

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left,
        top,
        width,
        thickness,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=SURFACE):
    """Add a rounded rectangle background shape."""
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x27, 0x27, 0x27)  # border color
    shape.line.width = Pt(1)
    # Adjust corner radius
    shape.adjustments[0] = 0.05  # subtle rounding
    return shape


def create_benefit_slide(prs, slide_data, slide_num):
    """Create a single key benefit slide with the Remaker Digital design system."""
    # Use blank layout so we have full control
    blank_layout = prs.slide_layouts[7]  # Layout 7: Blank
    slide = prs.slides.add_slide(blank_layout)

    # Slide dimensions
    sw = prs.slide_width
    sh = prs.slide_height

    # Margins
    margin_left = Inches(1.2)
    margin_right = Inches(1.2)
    content_width = sw - margin_left - margin_right

    # ── Background color (dark) ──
    # The template's master slide handles this via theme

    # ── Slide number indicator ──
    indicator_y = Inches(0.6)
    add_textbox(
        slide,
        margin_left,
        indicator_y,
        Inches(3),
        Inches(0.4),
        f"KEY BENEFIT {slide_num} OF 3",
        font_size=12,
        color=LIGHT_GRAY,
        bold=True,
        font_name="Montserrat",
    )

    # ── Brand accent line ──
    add_accent_line(slide, margin_left, Inches(1.1), Inches(2.5), color=BRAND_RED, thickness=Pt(4))

    # ── Main title ──
    add_textbox(
        slide,
        margin_left,
        Inches(1.4),
        content_width,
        Inches(1.2),
        slide_data["title"],
        font_size=44,
        color=WHITE,
        bold=True,
        font_name="Michroma",
    )

    # ── Subtitle ──
    add_textbox(
        slide,
        margin_left,
        Inches(2.6),
        content_width,
        Inches(0.6),
        slide_data["subtitle"],
        font_size=24,
        color=BRAND_RED,
        bold=True,
        font_name="Montserrat",
    )

    # ── Bullet points ──
    bullet_top = Inches(3.5)
    bullet_height = Inches(4.5)
    add_bullet_list(
        slide,
        margin_left,
        bullet_top,
        content_width,
        bullet_height,
        slide_data["points"],
        font_size=22,
        color=WHITE,
        bullet_color=BRAND_RED,
    )

    # ── Callout box at bottom ──
    callout_height = Inches(1.0)
    callout_top = sh - Inches(2.2)
    callout_left = margin_left
    callout_width = content_width

    # Rounded rectangle background for callout
    rect = add_rounded_rect(slide, callout_left, callout_top, callout_width, callout_height, fill_color=SURFACE)

    # Callout text
    add_textbox(
        slide,
        callout_left + Inches(0.4),
        callout_top + Inches(0.15),
        callout_width - Inches(0.8),
        callout_height - Inches(0.3),
        slide_data["callout"],
        font_size=20,
        color=BRAND_RED,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        font_name="Montserrat",
    )


def main():
    print(f"Loading template: {TEMPLATE}")
    prs = Presentation(str(TEMPLATE))

    # Change slide dimensions to 4:3 (1600x1200)
    print(f"Original dimensions: {prs.slide_width} x {prs.slide_height} EMU")
    prs.slide_width = TARGET_WIDTH
    prs.slide_height = TARGET_HEIGHT
    print(
        f'New dimensions: {prs.slide_width} x {prs.slide_height} EMU ({prs.slide_width / 914400:.2f}" x {prs.slide_height / 914400:.2f}")'
    )

    # Remove any existing slides from the template
    # Template may have placeholder slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        elem = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(elem)

    print(f"Cleared template slides. Creating {len(SLIDES)} key benefit slides...")

    # Create the 3 key benefit slides
    for i, slide_data in enumerate(SLIDES, 1):
        create_benefit_slide(prs, slide_data, i)
        print(f"  Created slide {i}: {slide_data['title']}")

    # Save
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f'Dimensions: {prs.slide_width / 914400:.2f}" x {prs.slide_height / 914400:.2f}" (4:3)')
    print(f"\nNext step: Open in PowerPoint and export each slide as PNG at 1600x1200")


if __name__ == "__main__":
    main()
